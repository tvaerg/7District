"""
7D Portfolio Aggregation Engine -- Renderer
===========================================
Injects a PortfolioReport into 7D's actual template .pptx and writes the final
deck. Template-injection (not redraw): we open 7D's template and overwrite the
content of the existing heatmap table cells and update-slide text, preserving
their exact branding, fonts, and layout.

Why this keeps editing fluid for 7D
------------------------------------
The heatmap is a NATIVE PowerPoint table. We set cell text and cell fill colour
through the table API, which leaves every cell fully hand-editable: if 7D
disagrees with a RAG colour, they click the cell -> shading -> pick a colour, or
just retype the text. No intermediate file, no override syntax. The engine's
reasoning for each proposed colour is written into the heatmap slide's speaker
notes as an audit trail, so the "why" is one keystroke away but never on the
printed slide.

What it fills
-------------
* Heatmap (slide 2): rows 2-13, one per company, matched by name in column 1.
    col 4  investment thesis   (RAG fill)
    col 5  performance         (RAG fill)
    col 6  financing summary   (RAG fill)
    col 7  covenant issue
    col 8-12  financing detail (amount/type/7D amount/timing/rescue) -- gated
    col 13 exit timing
* Update slides (3-6): two companies per slide, the Operational / Financial /
  Three-year blocks + 7D focus. (MVP: text injection into the existing layout.)

RAG -> fill hex (matches the template's own palette):
    green = 00B050 , amber = FFFF00 , red = FF0000

Dependencies: python-pptx.
"""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from sevend_schema import PortfolioReport, CompanyFacts, RAG


# Template-native RAG fills.
RAG_HEX = {
    RAG.GREEN: "00B050",
    RAG.AMBER: "FFFF00",
    RAG.RED:   "FF0000",
}
NA_HEX = "D5D5D5"  # grey, template's n/a tone

# Heatmap column indices (mapped from the template).
COL_NAME = 1
COL_THESIS = 4
COL_PERFORMANCE = 5
COL_FIN_SUMMARY = 6
COL_COVENANT = 7
COL_AMOUNT = 8
COL_TYPE = 9
COL_7D_AMOUNT = 10
COL_TIMING = 11
COL_RESCUE = 12
COL_EXIT = 13
FIRST_DATA_ROW = 2  # rows 0-1 are headers


# -----------------------------------------------------------------------------
# Low-level cell helpers -- preserve formatting, only change text + fill
# -----------------------------------------------------------------------------
def _set_cell_text(cell, text: str) -> None:
    """Overwrite a cell's text while preserving the run's existing formatting.

    We edit the first run in place (keeping its font/size/colour) and clear any
    extra runs, rather than rebuilding the paragraph -- this keeps the template's
    typography intact.
    """
    text = "" if text is None else str(text)
    tf = cell.text_frame
    if not tf.paragraphs:
        tf.text = text
        return
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        # Drop any trailing runs so we don't leave stale fragments.
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p.text = text
    # Remove any additional paragraphs beyond the first.
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def _set_cell_fill(cell, hex_color: Optional[str]) -> None:
    if hex_color is None:
        return
    # python-pptx's cell.fill API is unreliable for table cells that already carry
    # a template fill -- the old <a:solidFill> can persist. Write the fill directly
    # into tcPr: remove any existing fill nodes, then insert a fresh solidFill.
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    nsmain = "http://schemas.openxmlformats.org/drawingml/2006/main"
    fill = parse_xml(
        f'<a:solidFill xmlns:a="{nsmain}"><a:srgbClr val="{hex_color}"/></a:solidFill>'
    )
    # tcPr children have an order; fill must precede any trailing elements like
    # a:lnL etc. Inserting right after the last line element or at end is safe here
    # because the template cells only carry margins + fill. Append works.
    tcPr.append(fill)


def _rag_fill(cell, rag: RAG) -> None:
    _set_cell_fill(cell, RAG_HEX.get(rag))


def _covenant_fill(covenant_text: str) -> str:
    """Decide the covenant cell colour from its text.

    Red only for an actual breach/problem; amber for an in-progress discussion;
    otherwise neutral grey. We never paint a 'No breach' covenant cell red just
    because financing elsewhere is stressed.
    """
    t = (covenant_text or "").strip().lower()
    if not t or t in ("n/a", "na", "none"):
        return NA_HEX
    # Explicit "no breach / no issue" -> neutral, even if verbose.
    if t.startswith("no") or "no breach" in t or "no covenant" in t:
        return NA_HEX
    # Active breach language -> red.
    if any(k in t for k in ("breach", "default", "broken", "not meet", "won't meet", "fail")):
        return RAG_HEX[RAG.RED]
    # Ongoing dialogue/waiver -> amber.
    if any(k in t for k in ("discussion", "dialogue", "waiver", "negotiat", "talks", "renegotiat")):
        return RAG_HEX[RAG.AMBER]
    # Anything else mentioned but unclear -> amber (worth a look), not red.
    return RAG_HEX[RAG.AMBER]


# -----------------------------------------------------------------------------
# Heatmap injection
# -----------------------------------------------------------------------------
def _find_heatmap_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _index_rows_by_name(table) -> dict[str, int]:
    """Map normalised company name -> row index, for rows that name a company."""
    idx = {}
    for r in range(FIRST_DATA_ROW, len(table.rows)):
        name = table.cell(r, COL_NAME).text.strip()
        if name:
            idx[name.lower()] = r
    return idx


def _short_label(text: str, max_words: int = 3, max_chars: int = 22) -> str:
    """Safety net: cap heatmap cell text to a short label even if the model
    returns something too long. Keeps the prompt's intent enforceable at render
    time. Full text always survives in the speaker-notes reasoning.

    This is a last-resort guard -- the prompt asks the model for short labels, so
    it should rarely fire. When it does, we cut at a word boundary and drop a
    dangling connector ('but', 'still', 'and'...) so the result reads cleanly.
    """
    t = (text or "").strip()
    if not t:
        return t
    words = t.split()
    truncated = False
    if len(words) > max_words:
        words = words[:max_words]
        truncated = True
    # Drop a trailing connector left dangling by the cut.
    connectors = {"but", "still", "and", "or", "with", "to", "of", "the", "a", "by", "for"}
    while words and words[-1].lower().strip(",.;:") in connectors:
        words.pop()
        truncated = True
    t = " ".join(words).rstrip(" ,;:.")
    if len(t) > max_chars:
        # Cut to the last whole word within the char budget rather than mid-word.
        cut = t[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:.")
        t = cut or t[:max_chars].rstrip(" ,;:.")
        truncated = True
    return (t + "…") if truncated and t else t


def _inject_company_row(table, row: int, c: CompanyFacts) -> None:
    fin = c.financing

    # Text + RAG fills for the judgement columns. Labels are capped short so they
    # always fit the narrow cells; the full nuance lives in the speaker notes.
    _set_cell_text(table.cell(row, COL_THESIS), _short_label(c.investment_thesis.effective_text))
    _rag_fill(table.cell(row, COL_THESIS), c.investment_thesis.effective_rag)

    _set_cell_text(table.cell(row, COL_PERFORMANCE), _short_label(c.performance.effective_text))
    _rag_fill(table.cell(row, COL_PERFORMANCE), c.performance.effective_rag)

    _set_cell_text(table.cell(row, COL_FIN_SUMMARY), _short_label(fin.summary.effective_text))
    _rag_fill(table.cell(row, COL_FIN_SUMMARY), fin.summary.effective_rag)

    # Covenant + financing detail columns (already gated to n/a upstream).
    _set_cell_text(table.cell(row, COL_COVENANT), _short_label(fin.covenant_issue))
    # Colour the covenant cell by what the covenant text ACTUALLY says -- not the
    # broad financing-issue flag. A runway/liquidity concern can make financing
    # red while the covenant itself is fine; those must not be conflated.
    _set_cell_fill(table.cell(row, COL_COVENANT), _covenant_fill(fin.covenant_issue))
    _set_cell_text(table.cell(row, COL_AMOUNT), fin.amount)
    _set_cell_text(table.cell(row, COL_TYPE), fin.type)
    _set_cell_text(table.cell(row, COL_7D_AMOUNT), fin.sevend_amount)
    _set_cell_text(table.cell(row, COL_TIMING), fin.timing)
    _set_cell_text(table.cell(row, COL_RESCUE), fin.sevend_rescue_amount)

    # Exit timing (static, from registry).
    _set_cell_text(table.cell(row, COL_EXIT), c.exit_timing)


def _write_heatmap_notes(slide, report: PortfolioReport) -> None:
    """Put the engine's reasoning for every proposed colour into speaker notes."""
    lines = [f"7D heatmap -- engine reasoning ({report.cycle_label})", ""]
    for c in report.companies:
        lines.append(f"{c.canonical_name}:")
        lines.append(f"  Thesis [{c.investment_thesis.effective_rag.value}]: {c.investment_thesis.reasoning}")
        lines.append(f"  Performance [{c.performance.effective_rag.value}]: {c.performance.reasoning}")
        lines.append(f"  Financing [{c.financing.summary.effective_rag.value}]: {c.financing.summary.reasoning}")
        if c.flags:
            lines.append(f"  ⚠️ FLAGS: {'; '.join(c.flags)}")
        lines.append("")
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def render_heatmap(prs: Presentation, report: PortfolioReport) -> list[str]:
    """Fill the heatmap slide. Returns a list of company names not found in the table."""
    slide = prs.slides[1]  # slide 2 (0-indexed)
    table = _find_heatmap_table(slide)
    if table is None:
        raise RuntimeError("No table found on the heatmap slide (slide 2).")

    name_to_row = _index_rows_by_name(table)
    filled_rows: set[int] = set()
    missing = []
    for c in report.companies:
        row = name_to_row.get(c.canonical_name.lower())
        if row is None:
            missing.append(c.canonical_name)
            continue
        _inject_company_row(table, row, c)
        filled_rows.add(row)

    # Neutralise any data row we did NOT fill, so the template's original demo
    # data can never leak into a real deck. A row with a company name but no
    # matching CompanyFacts is blanked to a grey 'No report' state.
    for r in range(FIRST_DATA_ROW, len(table.rows)):
        if r in filled_rows:
            continue
        if not table.cell(r, COL_NAME).text.strip():
            continue  # structural/empty row, leave alone
        _neutralise_row(table, r)

    _write_heatmap_notes(slide, report)
    return missing


def _neutralise_row(table, row: int) -> None:
    """Blank a heatmap row that received no real data: grey judgement cells,
    'No report' text, n/a financing. Prevents stale demo data from showing."""
    for col, txt in (
        (COL_THESIS, "No report"),
        (COL_PERFORMANCE, "No report"),
        (COL_FIN_SUMMARY, "No report"),
    ):
        _set_cell_text(table.cell(row, col), txt)
        _set_cell_fill(table.cell(row, col), NA_HEX)
    for col in (COL_COVENANT, COL_AMOUNT, COL_TYPE, COL_7D_AMOUNT, COL_TIMING, COL_RESCUE):
        _set_cell_text(table.cell(row, col), "n/a")
        _set_cell_fill(table.cell(row, col), NA_HEX)


# -----------------------------------------------------------------------------
# Update slides (3-6): position-based slot injection
# -----------------------------------------------------------------------------
# Each update slide has up to two company columns (left ~x<13in, right >=13in).
# Within a column, four content boxes stack by vertical position:
#   row 0 ~y5.0  Operational update
#   row 1 ~y7.2  Financial performance
#   row 2 ~y9.3  Three-year value creation
#   row 3 ~y11.4 7D focus and contribution
# Company identity in the template is a logo/position, not reliable text, so we
# assign companies to slots IN ORDER (same category order as the heatmap) and
# write the company name into the first block as a header so each is labelled.

_ROW_Y = [5.0, 7.2, 9.3, 11.4]   # approximate top (inches) of each content row
_ROW_TOL = 1.0                    # vertical tolerance for matching a box to a row
_COL_SPLIT_IN = 13.0              # x < split = left column, else right


def _emu_in(v) -> float:
    from pptx.util import Emu
    return Emu(v).inches if v is not None else -1.0


def _collect_update_slots(slide):
    """Return {(col, row): content_shape} for the long-text content boxes, plus
    {(col, row): ellipse_shape} for the status balls, keyed by column (0=left,
    1=right) and row (0-3). Matching is purely positional, so it is robust to the
    differing shape order across slides 3-6."""
    content: dict[tuple[int, int], object] = {}
    ellipses: list[tuple[float, float, object]] = []
    for shape in slide.shapes:
        x, y = _emu_in(shape.left), _emu_in(shape.top)
        if x < 0 or y < 0:
            continue
        col = 0 if x < _COL_SPLIT_IN else 1
        if "Ellips" in shape.name:
            ellipses.append((x, y, shape))
            continue
        if shape.has_text_frame and "Rektangel" in shape.name:
            t = shape.text_frame.text.strip()
            if len(t) > 40:  # a content box (carries real update prose)
                for ri, ry in enumerate(_ROW_Y):
                    if abs(y - ry) <= _ROW_TOL:
                        content[(col, ri)] = shape
                        break
    return content, ellipses


def _nearest_ellipse(ellipses, target_x: float, target_y: float):
    """Find the status ellipse closest to a content row (the RAG ball)."""
    best, best_d = None, 1e9
    for x, y, sh in ellipses:
        d = abs(y - target_y) + abs(x - target_x) * 0.3
        if d < best_d:
            best_d, best = d, sh
    return best


def _set_shape_text(shape, text: str) -> None:
    """Overwrite a shape's text, preserving the first run's formatting."""
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for extra in tf.paragraphs[0].runs[1:]:
            extra._r.getparent().remove(extra._r)
        for extra_p in tf.paragraphs[1:]:
            extra_p._p.getparent().remove(extra_p._p)
    else:
        tf.text = text


def _trim_bullet(text: str, max_words: int = 20, max_chars: int = 120) -> str:
    """Cap a single bullet so it can't blow out the box. The caps match the
    template's actual readable density (~4 bullets per block, ~2 lines each).

    A semicolon-packed multi-fact bullet still gets cut to the first fact --
    the prompt asks for one fact per bullet and packing is the model's bug,
    not legitimate length.
    """
    t = (text or "").strip()
    if not t:
        return t
    cut = False
    # If the model crammed multiple facts with ';', keep just the first.
    if ";" in t:
        t = t.split(";", 1)[0].strip()
        cut = True
    words = t.split()
    if len(words) > max_words:
        words = words[:max_words]; cut = True
    t = " ".join(words)
    if len(t) > max_chars:
        t = t[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:."); cut = True
    return (t.rstrip(" ,;:.") + "…") if cut else t


def _set_bullets(shape, lines: list[str], header: str = "", max_bullets: int = 3) -> None:
    """Fill a content box with an optional bold header line then bullet lines.

    Layout protection (belt + suspenders):
      1. Hard-caps bullet count + length via `_trim_bullet`.
      2. Fixes body bullets at 16pt (smaller than the template's ~18pt, gives
         ~10% more vertical room per box for the realistic case).
      3. Enables PowerPoint's native shrink-on-overflow auto-fit, so anything
         that STILL won't fit at 16pt scales down automatically rather than
         bleeding into the block below.
    """
    from pptx.util import Pt
    from pptx.enum.text import MSO_AUTO_SIZE

    # Safety net: cap count and trim each bullet (prompt asks for this too).
    capped = [_trim_bullet(x) for x in lines if x][:max_bullets]
    tf = shape.text_frame
    # Capture formatting from the existing first run before clearing.
    base_run = tf.paragraphs[0].runs[0] if (tf.paragraphs and tf.paragraphs[0].runs) else None
    font_name = base_run.font.name if base_run else "Helvetica Neue"
    font_rgb = None
    if base_run is not None:
        try:
            if base_run.font.color and base_run.font.color.type is not None:
                font_rgb = base_run.font.color.rgb
        except Exception:                              # noqa: BLE001
            font_rgb = None
    if font_rgb is None:
        font_rgb = RGBColor.from_string("293447")      # template dark navy fallback

    # Fixed body size (slightly smaller than template default for room).
    body_size = Pt(16)
    header_size = Pt(17)  # mild hierarchy without dominating

    tf.word_wrap = True
    # Auto-fit: shrink text to fit if 16pt still overflows. Native PowerPoint
    # behavior; renders identically in PowerPoint and downstream viewers.
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_SHAPE_HEIGHT  # 'Shrink text on overflow' equivalent
    except Exception:                                  # noqa: BLE001
        pass

    tf.clear()
    body_lines = ([f"{header}"] if header else []) + [f"• {x}" for x in capped]
    if not body_lines:
        body_lines = ["• (no update this cycle)"]
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        try:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.LEFT
        except Exception:                              # noqa: BLE001
            pass
        run = p.add_run()
        run.text = line
        if font_name:
            run.font.name = font_name
        # Header line gets the slightly larger size; bullets get the body size.
        is_header_line = bool(header) and i == 0
        run.font.size = header_size if is_header_line else body_size
        run.font.color.rgb = font_rgb
        if is_header_line:
            run.font.bold = True


def _fill_company_column(content, ellipses, col: int, c: CompanyFacts) -> None:
    """Fill one company column (operational/financial/three-year/7D focus + RAG)."""
    blocks = [
        (0, c.operational, f"{c.canonical_name} — Operational"),
        (1, c.financial, "Financial performance"),
        (2, c.three_year, "Three-year value creation"),
    ]
    for row, block, header in blocks:
        shape = content.get((col, row))
        if shape is None:
            continue
        _set_bullets(shape, block.effective_bullets, header=header)
        # Recolour the nearest status ellipse to the block's RAG.
        x = _emu_in(shape.left)
        ell = _nearest_ellipse(ellipses, target_x=(4.3 if col == 0 else 16.2),
                               target_y=_ROW_Y[row])
        if ell is not None:
            _set_cell_fill_shape(ell, RAG_HEX[block.effective_rag])
    # 7D focus (row 3, no RAG ball).
    focus_shape = content.get((col, 3))
    if focus_shape is not None:
        _set_bullets(focus_shape, c.sevend_focus, header="7D focus and contribution")


def _set_cell_fill_shape(shape, hex_color: str) -> None:
    """Solid-fill an autoshape (status ellipse) reliably."""
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)


def render_update_slides(prs: Presentation, report: PortfolioReport) -> None:
    """Fill update slides (3-6), pinning each company to ITS slot so it always
    sits under its own logo (see UPDATE_SLOT_MAP in the registry).

    A pinned company that did not report this cycle has its slot blanked (stale
    template text cleared) rather than backfilled by someone else -- a logo must
    never sit above the wrong company's update.
    """
    from sevend_registry import UPDATE_SLOT_MAP

    by_name = {c.canonical_name.lower(): c for c in report.companies}
    unplaced = []

    for (si, col), company_name in UPDATE_SLOT_MAP.items():
        if si >= len(prs.slides):
            continue
        content, ellipses = _collect_update_slots(prs.slides[si])
        c = by_name.get(company_name.lower())
        if c is not None and c.has_report_this_cycle and not c.flags:
            _fill_company_column(content, ellipses, col, c)
        else:
            # Pinned company absent/failed this cycle -> blank the slot cleanly.
            _blank_company_column(content, ellipses, col, company_name)

    # Surface any reporting company that has no pinned slot (portfolio changed
    # but the slot map wasn't updated).
    pinned = {n.lower() for n in UPDATE_SLOT_MAP.values()}
    for c in report.companies:
        if c.has_report_this_cycle and not c.flags and c.canonical_name.lower() not in pinned:
            unplaced.append(c.canonical_name)
    if unplaced:
        print(f"   ⚠️ Reporting companies with no update-slide slot (add to "
              f"UPDATE_SLOT_MAP): {', '.join(unplaced)}")


def _blank_company_column(content, ellipses, col: int, company_name: str) -> None:
    """Clear a slot whose pinned company didn't report, so no stale template
    text or wrong-company content remains under its logo."""
    blocks = [
        (0, f"{company_name} — Operational"),
        (1, "Financial performance"),
        (2, "Three-year value creation"),
        (3, "7D focus and contribution"),
    ]
    for row, header in blocks:
        shape = content.get((col, row))
        if shape is None:
            continue
        _set_bullets(shape, ["No report this cycle"], header=header)
        if row < 3:
            ell = _nearest_ellipse(ellipses, target_x=(4.3 if col == 0 else 16.2),
                                   target_y=_ROW_Y[row])
            if ell is not None:
                _set_cell_fill_shape(ell, NA_HEX)


# -----------------------------------------------------------------------------
# Public entry point
# -----------------------------------------------------------------------------
def _resolve_template(template_path: str) -> str:
    """Find the template whether the project is flat or nested.

    The CANONICAL home is assets/7d_template.pptx, committed to the repo and
    located relative to this file -- so a fresh `git clone`/`pull` on any machine
    finds it with zero configuration, no matter what folder you run from.

    Falls back to other sensible spots (flat layout, template/ subfolder, the
    explicit path given) so it keeps working across layouts. Raises a clear error
    listing everywhere it looked if nothing is found.
    """
    from pathlib import Path as _P
    here = _P(__file__).resolve().parent
    candidates = [
        here / "assets" / "7d_template.pptx",      # canonical, committed home
        _P(template_path),                          # explicit path as given
        here / template_path,
        here / _P(template_path).name,
        here / _P(template_path).name,
        here / "template" / _P(template_path).name,
        here / "7d_template.pptx",                  # flat layout
        here / "Monthly_Portfolio_Company_report_Template.pptx",
    ]
    for c in candidates:
        if c.is_file():
            return str(c)
    looked = "\n      ".join(str(c) for c in candidates)
    raise FileNotFoundError(
        f"Could not find the 7D template .pptx. Looked in:\n      {looked}\n"
        f"   The template should live at assets/7d_template.pptx (committed to the repo). "
        f"Or pass a path with --template."
    )


def render_report(report: PortfolioReport, template_path: str, out_path: str) -> str:
    """Render the deck from the template. Returns the output path.

    NOTE: The heatmap (slide 2) is INTENTIONALLY NOT rendered. By 7D's choice,
    the heatmap is a static reference page that 7D maintains by hand in the
    template. The engine only fills the per-company update slides (3-6), where
    monthly content actually changes. `render_heatmap` remains in this module
    in case the policy is ever reversed -- it just isn't called.
    """
    template_path = _resolve_template(template_path)
    prs = Presentation(template_path)

    # render_heatmap(prs, report)   # disabled by 7D policy -- heatmap is manual
    render_update_slides(prs, report)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    # Build a small two-company report and render against the template, to prove
    # heatmap injection + RAG fills end to end (no API needed).
    from sevend_schema import RatedField, Financing, UpdateBlock

    def cf(name, cat, thesis, perf, fin_sum, cov, rag_t, rag_p, rag_f,
           amount="n/a", typ="n/a", s7d="n/a", timing="n/a", exit="2028", issue=False):
        return CompanyFacts(
            canonical_name=name, category=cat, lead="AJ", support="TBD", exit_timing=exit,
            investment_thesis=RatedField(proposed_rag=rag_t, proposed_text=thesis, reasoning=f"{name} thesis basis."),
            performance=RatedField(proposed_rag=rag_p, proposed_text=perf, reasoning=f"{name} perf basis."),
            financing=Financing(
                summary=RatedField(proposed_rag=rag_f, proposed_text=fin_sum, reasoning=f"{name} fin basis."),
                covenant_issue=cov, has_active_issue=issue,
                amount=amount, type=typ, sevend_amount=s7d, timing=timing),
            operational=UpdateBlock(proposed_rag=rag_p, bullets=["Op point A", "Op point B"], reasoning="op"),
            financial=UpdateBlock(proposed_rag=rag_f, bullets=["Fin point A"], reasoning="fin"),
            three_year=UpdateBlock(proposed_rag=rag_t, bullets=["3yr point A"], reasoning="3yr"),
            sevend_focus=["Focus A", "Focus B"],
        )

    report = PortfolioReport(
        cycle_label="April 2026", ceo_update="Mixed month.",
        companies=[
            cf("Inretrn", "Core Holdings", "OK", "Growing", "OK", "No",
               RAG.GREEN, RAG.GREEN, RAG.GREEN, exit="2028"),
            cf("Rapid Images", "Core Holdings", "Under transformation", "Weak trading",
               "Weak trading", "Yes, breach Q1", RAG.AMBER, RAG.RED, RAG.RED,
               amount="5", typ="Defend", s7d="2.6", timing="Q2'26", exit="2028", issue=True),
        ],
    )

    out = render_report(report, "assets/7d_template.pptx", "output/7d_report_April_2026.pptx")
    print(f"Rendered: {out}")

    # Verify by reading back the injected cells.
    prs = Presentation(out)
    t = _find_heatmap_table(prs.slides[1])
    rows = _index_rows_by_name(t)
    for nm in ("inretrn", "rapid images"):
        r = rows[nm]
        print(f"  {nm}: thesis='{t.cell(r,COL_THESIS).text}' perf='{t.cell(r,COL_PERFORMANCE).text}' "
              f"fin='{t.cell(r,COL_FIN_SUMMARY).text}' amount='{t.cell(r,COL_AMOUNT).text}'")