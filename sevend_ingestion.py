"""
7D Portfolio Aggregation Engine -- Ingestion Layer
==================================================
Turns a heterogeneous pile of monthly portfolio-company reports into a uniform
intermediate form (RawCompanyMaterial) that the reasoning layer can consume,
regardless of the original file format or the company's industry/KPI language.

Two format paths (per the agreed design):

    PDF   -> Gemini native PDF read. Gemini ingests the file directly (text,
             tables, charts) via the File API. No rasterisation. A two-phase
             CUSTODIAL read runs: (1) a holistic "owner's eye" deep read that
             captures soft signal not in the KPI tables, then (2) a faithful
             transcription for the reasoning layer.

    PPTX  -> Defensive text + table + chart + notes extraction (python-pptx),
             walking nested group shapes. PPTX is already structured text, so
             it goes straight to Opus in the reasoning layer; no vision pass.

Robustness (why this is more than a happy-path reader)
------------------------------------------------------
* Retry + backoff on EVERY Gemini call (upload AND generate), not just upload.
* PDF size / page-count guard with a warning before burning tokens.
* Defensive PPTX walk: recurses into grouped shapes, pulls chart category/series
  data, and captures speaker notes -- so data buried off the main text frame is
  not silently lost.
* Every failure is captured on the material's `errors` list and surfaced loudly;
  a company is never dropped silently.

Design notes
------------
* SELF-CONTAINED. No imports from VAERG. Patterns are borrowed (Gemini upload +
  patterns are borrowed (Gemini upload + poll + cleanup; robust retry), the
  code is rewritten for handover to 7D.
* CUSTODIAL lens. These are 7D's OWN portfolio companies reporting to their
  owner -- not a seller pitching an IM. The deep read looks for what the owner
  must ACT ON, not for what someone is "hiding".
* This layer does NOT make RAG-colour judgements and does NOT fill the 7D
  schema. It produces clean, labelled material + soft-signal notes only.

Env:
    GEMINI_API_KEY   -- required for the PDF path.
"""

from __future__ import annotations

import os
import time
import json
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional, Callable, Any

from dotenv import load_dotenv

from sevend_registry import match_file_to_company, CompanyRecord

# Load .env robustly regardless of project layout: try the same dir as this
# file (flat layout: code + .env together) AND the parent dir (src/ layout),
# falling back to the default CWD search. First hit wins.
def _load_env():
    here = Path(__file__).resolve().parent
    for candidate in (here / ".env", here.parent / ".env"):
        if candidate.is_file():
            load_dotenv(candidate)
            return
    load_dotenv()  # fall back to default search (CWD and upward)

_load_env()

# --- Configuration -----------------------------------------------------------
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "models/gemini-3.1-pro-preview")  # native PDF
_PDF_EXTS = {".pdf"}
_PPTX_EXTS = {".pptx", ".ppt"}
_XLSX_EXTS = {".xlsx", ".xlsm"}
_DOCX_EXTS = {".docx", ".doc"}

# Guard rails. Reports are normally 5-25 pages; flag anything unusually large
# before we spend tokens, but do not hard-block (operator may legitimately have
# a big appendix-heavy report).
MAX_PDF_MB = 40
WARN_PDF_PAGES = 60

# Retry policy for Gemini calls.
_MAX_RETRIES = 5
_BACKOFF_BASE_S = 3  # 3, 6, 9, 12 ...


# -----------------------------------------------------------------------------
# Intermediate form -- the contract between ingestion and reasoning
# -----------------------------------------------------------------------------
@dataclass
class RawCompanyMaterial:
    """Everything the reasoning layer needs about ONE company, format-agnostic.

    `body_text` is the full report content as clean text (Gemini's faithful
    transcription for PDFs; defensive python-pptx extraction for PPTX).
    `deep_read` holds the custodial soft-signal notes (PDF path). `source_kind`
    tells the reasoning layer how much to trust table fidelity. `warnings`
    holds non-fatal issues (e.g. oversized PDF) distinct from
    `errors` (which mean the material is incomplete or unusable).
    """
    canonical_name: str
    category: str
    source_file: str
    source_kind: str                       # "pdf_gemini" | "pptx_text" | "xlsx_kpi" | "unknown"
    body_text: str = ""
    deep_read: Optional[dict] = None        # custodial notes; None for pptx/xlsx
    static: dict = field(default_factory=dict)   # registry static fields
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    meta: dict = field(default_factory=dict)     # page count, bytes, etc.

    # Multi-file bundle support: when several files exist for one company they
    # are MERGED into one material (see bundle_materials() in the pipeline).
    # `sources` lists every file that contributed; `kpis` carries deterministic
    # numeric snapshots from any xlsx files.
    sources: list[str] = field(default_factory=list)
    kpis: list[dict] = field(default_factory=list)  # one KpiSnapshot.to_json() per xlsx

    @property
    def ok(self) -> bool:
        return not self.errors and bool(self.body_text)

    def to_json(self) -> dict:
        return asdict(self)


# -----------------------------------------------------------------------------
# Shared helpers
# -----------------------------------------------------------------------------
def _robust_json(text: str) -> Optional[dict]:
    """Parse JSON that may be wrapped in prose or code fences."""
    if not text:
        return None
    t = text.strip()
    if t.startswith("```"):
        # Strip a leading ```json / ``` fence and any trailing fence.
        t = t[3:]
        if t.lower().startswith("json"):
            t = t[4:]
        if "```" in t:
            t = t[: t.rfind("```")]
        t = t.strip()
    try:
        return json.loads(t)
    except Exception:                                  # noqa: BLE001
        pass
    start, end = t.find("{"), t.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(t[start : end + 1])
        except Exception:                              # noqa: BLE001
            return None
    return None


def _with_retries(fn: Callable[[], Any], label: str) -> Any:
    """Run a Gemini call with exponential backoff. Raises on final failure."""
    last = None
    for attempt in range(1, _MAX_RETRIES + 1):
        try:
            return fn()
        except Exception as err:                       # noqa: BLE001
            last = err
            if attempt < _MAX_RETRIES:
                wait = _BACKOFF_BASE_S * attempt
                print(f"      ⚠️ {label} attempt {attempt}/{_MAX_RETRIES} failed: "
                      f"{str(err)[:80]} -- retry in {wait}s")
                time.sleep(wait)
    raise RuntimeError(f"{label} failed after {_MAX_RETRIES} attempts: {last}")


# -----------------------------------------------------------------------------
# PDF PATH -- Gemini native, two-phase custodial
# -----------------------------------------------------------------------------
def _configure_gemini():
    try:
        import google.generativeai as genai
    except ImportError:
        raise RuntimeError(
            "google-generativeai not installed -- run: pip install -r requirements.txt"
        )
    key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    if not key:
        raise RuntimeError(
            "Neither GEMINI_API_KEY nor GOOGLE_API_KEY set -- required for the PDF path."
        )
    genai.configure(api_key=key)
    return genai


def _pdf_page_count(pdf_path: str) -> Optional[int]:
    """Best-effort page count. Tries pypdf, then a light byte scan. None if unknown."""
    try:
        from pypdf import PdfReader
        return len(PdfReader(pdf_path).pages)
    except Exception:                                  # noqa: BLE001
        pass
    try:
        data = Path(pdf_path).read_bytes()
        # Rough heuristic: count /Type /Page objects (not /Pages).
        import re
        n = len(re.findall(rb"/Type\s*/Page[^s]", data))
        return n or None
    except Exception:                                  # noqa: BLE001
        return None


def _gemini_upload(genai, pdf_path: str):
    """Upload + poll until ready, with retry/backoff."""
    def _do_upload():
        return genai.upload_file(
            path=pdf_path, display_name=Path(pdf_path).stem
        )
    uploaded = _with_retries(_do_upload, "PDF upload")

    # Poll for processing with a sane ceiling so we never hang forever.
    waited = 0
    while uploaded.state.name == "PROCESSING":
        time.sleep(2)
        waited += 2
        uploaded = genai.get_file(uploaded.name)
        if waited > 180:
            raise RuntimeError("Gemini processing exceeded 180s -- aborting.")
    if uploaded.state.name == "FAILED":
        raise RuntimeError("Gemini failed to process the PDF.")
    return uploaded


_DEEP_READ_PROMPT = """You are a portfolio director at 7D (a holding company) reading the latest \
monthly report from a company you OWN. You are not a sceptical buyer dissecting a sales pitch -- \
you are the owner, so you need the TRUTH, including the uncomfortable parts, because it is your \
capital and you may have to act on it.

Read the ENTIRE report holistically first -- text, tables, AND charts. Charts and tables contain \
quantified data; read them with the same care as prose.

Your goal is NOT to extract every number yet. It is to capture the SOFT SIGNAL that a plain table \
of KPIs would miss. Produce notes on exactly four things:

1. NARRATIVE & TONE
   - What is management's story this month, in one or two sentences?
   - How confident or worried does the writing sound? Quote a short phrase if it is telling.

2. SOFT RISK NOT IN THE KPI TABLE
   - Anything that lives in the prose, not the numbers: covenant or bank discussions, refinancing,
     liquidity/runway concerns, key-person departures, customer dependency, factory/operational
     problems, regulatory or supply issues. List each as a short, specific bullet.

3. NARRATIVE-VS-NUMBERS GAPS
   - Where does the upbeat language not match the figures (e.g. "strong month" while EBITDA is
     below budget, or growth claimed while a key segment shrinks)? Flag each gap explicitly.

4. FINANCING & EXIT SIGNALS
   - Anything bearing on capital need, covenant status, debt, share issues, or exit/sale timing.

Return ONLY valid JSON, no markdown fences, in this exact shape:
{
  "narrative_and_tone": {"summary": "", "confidence_read": "", "telling_phrase": ""},
  "soft_risks": ["", ""],
  "narrative_vs_numbers_gaps": ["", ""],
  "financing_exit_signals": ["", ""]
}
If a list has nothing to report, return an empty list. Do not invent items."""

_TRANSCRIBE_PROMPT = """Transcribe this portfolio-company monthly report into clean, faithful plain \
text for downstream analysis. Preserve ALL quantified data.

Rules:
- Render every table as readable text, keeping row/column relationships clear (label: value, or
  aligned rows). Do not drop numbers.
- Convert charts and infographics into their underlying data points where legible (e.g. "MRR Mar
  2026: 2 132 903"). Charts ARE data.
- Keep section headings and the order of the document.
- Do not summarise, editorialise, or omit. This is a transcription, not an analysis.
- Output plain text only."""


def ingest_pdf(pdf_path: str, rec: CompanyRecord) -> RawCompanyMaterial:
    """PDF path: Gemini native read -> transcription + custodial deep-read notes."""
    material = RawCompanyMaterial(
        canonical_name=rec.canonical_name, category=rec.category,
        source_file=str(pdf_path), source_kind="pdf_gemini",
        static=rec.display_dict(),
    )

    # --- Pre-flight checks (cheap, before any API spend) ---
    p = Path(pdf_path)
    if not p.exists():
        material.errors.append(f"file not found: {pdf_path}")
        return material
    size_mb = p.stat().st_size / (1024 * 1024)
    material.meta["size_mb"] = round(size_mb, 2)
    if size_mb == 0:
        material.errors.append("file is empty (0 bytes)")
        return material
    if size_mb > MAX_PDF_MB:
        material.warnings.append(
            f"PDF is {size_mb:.0f}MB (> {MAX_PDF_MB}MB guard) -- proceeding but watch cost."
        )
    pages = _pdf_page_count(pdf_path)
    if pages is not None:
        material.meta["pages"] = pages
        if pages > WARN_PDF_PAGES:
            material.warnings.append(
                f"PDF has {pages} pages (> {WARN_PDF_PAGES}) -- large context, watch cost/limits."
            )

    # --- Gemini calls (always fresh; no caching) ---
    try:
        genai = _configure_gemini()
    except RuntimeError as e:
        material.errors.append(str(e))
        return material

    uploaded = None
    try:
        print(f"      📄 uploading to Gemini...")
        uploaded = _gemini_upload(genai, pdf_path)
        model = genai.GenerativeModel(GEMINI_MODEL)

        print(f"      🔎 phase 1: custodial deep read...")
        dr = _with_retries(
            lambda: model.generate_content([uploaded, _DEEP_READ_PROMPT]),
            "deep read",
        )
        parsed = _robust_json(getattr(dr, "text", "") or "")
        if parsed is None:
            material.warnings.append("deep-read JSON parse failed -- kept raw text")
            material.deep_read = {"_raw": getattr(dr, "text", "")}
        else:
            material.deep_read = parsed

        print(f"      📝 phase 2: faithful transcription...")
        tr = _with_retries(
            lambda: model.generate_content([uploaded, _TRANSCRIBE_PROMPT]),
            "transcription",
        )
        material.body_text = (getattr(tr, "text", "") or "").strip()
        if not material.body_text:
            material.errors.append("Gemini returned an empty transcription")

    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"gemini_error: {e}")
    finally:
        if uploaded is not None:
            try:
                genai.delete_file(uploaded.name)
            except Exception:                          # noqa: BLE001
                pass
    return material


# -----------------------------------------------------------------------------
# PPTX PATH -- defensive text + table + chart + notes, recursing groups
# -----------------------------------------------------------------------------
def _walk_shapes(shapes, lines: list[str]) -> None:
    """Recursively pull text, tables, and chart data from a shape tree.

    python-pptx hides content in places a naive `shape.text` misses: grouped
    shapes (need recursion), tables (cell-by-cell), and charts (category/series
    data lives on the chart's plots, not in any text frame). We pull all three.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        # Grouped shapes: recurse.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(shape.shapes, lines)
            continue

        # Text frames.
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt:
                    lines.append(txt)

        # Tables.
        if shape.has_table:
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    lines.append(" | ".join(cells))

        # Charts -- extract category labels and series values as data points.
        if getattr(shape, "has_chart", False):
            try:
                chart = shape.chart
                cats = [str(c) for c in chart.plots[0].categories] if chart.plots else []
                lines.append("[chart data]")
                for series in chart.series:
                    vals = list(series.values)
                    name = series.name or "series"
                    if cats and len(cats) == len(vals):
                        pairs = ", ".join(f"{c}: {v}" for c, v in zip(cats, vals))
                        lines.append(f"  {name}: {pairs}")
                    else:
                        lines.append(f"  {name}: {vals}")
            except Exception:                          # noqa: BLE001
                lines.append("[chart present but data not machine-readable]")


def ingest_pptx(pptx_path: str, rec: CompanyRecord) -> RawCompanyMaterial:
    """PPTX path: defensive extraction with python-pptx. No vision pass."""
    material = RawCompanyMaterial(
        canonical_name=rec.canonical_name, category=rec.category,
        source_file=str(pptx_path), source_kind="pptx_text",
        static=rec.display_dict(),
    )
    p = Path(pptx_path)
    if not p.exists():
        material.errors.append(f"file not found: {pptx_path}")
        return material
    if p.stat().st_size == 0:
        material.errors.append("file is empty (0 bytes)")
        return material

    try:
        from pptx import Presentation
    except ImportError:
        material.errors.append("python-pptx not installed (pip install python-pptx)")
        return material

    try:
        prs = Presentation(pptx_path)
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"could not open pptx (corrupt or not a pptx?): {e}")
        return material

    try:
        lines: list[str] = []
        slide_count = 0
        for idx, slide in enumerate(prs.slides, 1):
            slide_count += 1
            lines.append(f"\n## Slide {idx}")
            _walk_shapes(slide.shapes, lines)
            # Speaker notes often carry context not on the slide.
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    lines.append(f"[notes] {notes}")
        material.meta["slides"] = slide_count
        material.body_text = "\n".join(lines).strip()
        if not material.body_text:
            material.errors.append("pptx produced no extractable text")
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"pptx_extraction_error: {e}")
    return material


# -----------------------------------------------------------------------------
# XLSX PATH -- deterministic KPI extraction + readable text dump for context
# -----------------------------------------------------------------------------
def ingest_xlsx(xlsx_path: str, rec: CompanyRecord) -> RawCompanyMaterial:
    """Excel path: pull structured KPIs deterministically, and dump readable
    sheet text as backup context. Both go into RawCompanyMaterial."""
    material = RawCompanyMaterial(
        canonical_name=rec.canonical_name, category=rec.category,
        source_file=str(xlsx_path), source_kind="xlsx_kpi",
        static=rec.display_dict(),
    )
    p = Path(xlsx_path)
    if not p.exists():
        material.errors.append(f"file not found: {xlsx_path}")
        return material
    if p.stat().st_size == 0:
        material.errors.append("file is empty (0 bytes)")
        return material

    # Deterministic KPI extraction.
    from sevend_xlsx_kpi import extract_kpis_from_xlsx, render_kpis_for_prompt
    try:
        snap = extract_kpis_from_xlsx(xlsx_path)
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"kpi extractor crashed: {e}")
        return material
    if not snap.is_empty():
        material.kpis.append(snap.to_json())
    else:
        material.warnings.append(
            "no KPIs extracted (no recognised labels); xlsx will only contribute raw text"
        )
        if snap.notes:
            material.warnings.append("; ".join(snap.notes))

    # Readable text dump of every sheet as backup context for Opus.
    try:
        from openpyxl import load_workbook
        wb = load_workbook(p, data_only=True, read_only=True)
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"could not open xlsx: {e}")
        return material

    try:
        lines: list[str] = []
        for sn in wb.sheetnames:
            ws = wb[sn]
            lines.append(f"\n## Sheet: {sn}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                if any(c for c in cells):
                    lines.append(" | ".join(cells))
        material.body_text = "\n".join(lines).strip()
        # Prepend the KPI snapshot so it's the first thing Opus reads.
        kpi_block = render_kpis_for_prompt(snap)
        if kpi_block:
            material.body_text = kpi_block + "\n\n" + material.body_text
        material.meta["sheets"] = wb.sheetnames
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"xlsx_extraction_error: {e}")
    return material


# -----------------------------------------------------------------------------
# DOCX PATH -- python-docx text extraction, no API
# -----------------------------------------------------------------------------
def ingest_docx(docx_path: str, rec: CompanyRecord) -> RawCompanyMaterial:
    """Word document path: walk paragraphs + tables, return text. No LLM."""
    material = RawCompanyMaterial(
        canonical_name=rec.canonical_name, category=rec.category,
        source_file=str(docx_path), source_kind="docx_text",
        static=rec.display_dict(),
    )
    p = Path(docx_path)
    if not p.exists():
        material.errors.append(f"file not found: {docx_path}")
        return material
    if p.stat().st_size == 0:
        material.errors.append("file is empty (0 bytes)")
        return material

    try:
        import docx
    except ImportError:
        material.errors.append("python-docx not installed (pip install python-docx)")
        return material

    try:
        doc = docx.Document(str(p))
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"could not open docx: {e}")
        return material

    try:
        chunks: list[str] = []
        for para in doc.paragraphs:
            t = (para.text or "").strip()
            if t:
                chunks.append(t)
        # Tables can carry the financial detail in management reports.
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    chunks.append(" | ".join(cells))
        material.body_text = "\n".join(chunks).strip()
        material.meta["paragraphs"] = len(doc.paragraphs)
        material.meta["tables"] = len(doc.tables)
        if not material.body_text:
            material.warnings.append("docx parsed but no text found")
    except Exception as e:                             # noqa: BLE001
        material.errors.append(f"docx extraction error: {e}")
    return material


# -----------------------------------------------------------------------------
# Dispatcher
# -----------------------------------------------------------------------------
def ingest_file(file_path: str) -> RawCompanyMaterial:
    """Resolve company by reading file CONTENT (not filename), then route by
    extension to the right ingestion path."""
    from sevend_match import identify_company
    rec, why = identify_company(file_path)
    if rec is None:
        m = RawCompanyMaterial(
            canonical_name="UNMATCHED", category="UNMATCHED",
            source_file=str(file_path), source_kind="unknown",
        )
        m.errors.append(
            f"Could not identify '{Path(file_path).name}' as a registry company "
            f"(reason: {why}). Add an alias to sevend_registry.py, set up a "
            f"FILE_MAP override, or ensure the company name appears in the file."
        )
        return m

    ext = Path(file_path).suffix.lower()
    print(f"   • {rec.canonical_name} ({Path(file_path).name})  [{why}]")
    if ext in _PDF_EXTS:
        return ingest_pdf(file_path, rec)
    if ext in _PPTX_EXTS:
        return ingest_pptx(file_path, rec)
    if ext in _XLSX_EXTS:
        return ingest_xlsx(file_path, rec)
    if ext in _DOCX_EXTS:
        return ingest_docx(file_path, rec)

    m = RawCompanyMaterial(
        canonical_name=rec.canonical_name, category=rec.category,
        source_file=str(file_path), source_kind="unknown", static=rec.display_dict(),
    )
    m.errors.append(f"Unsupported file type: {ext} (expected .pdf, .pptx, .xlsx, or .docx)")
    return m


def ingest_portfolio(file_paths: list[str]) -> list[RawCompanyMaterial]:
    """Ingest all monthly reports for this cycle into uniform material.

    Surfaces both hard errors and soft warnings, and reports which registry
    companies received NO report this cycle (so the heatmap will have gaps).
    """
    print("📥 INGESTION")
    out: list[RawCompanyMaterial] = []
    for fp in file_paths:
        out.append(ingest_file(fp))

    # Coverage check: which registry companies got no file this cycle?
    from sevend_registry import REGISTRY
    covered = {m.canonical_name for m in out if m.canonical_name != "UNMATCHED"}
    missing = [r.canonical_name for r in REGISTRY if r.canonical_name not in covered]

    errs = [m for m in out if m.errors]
    warns = [m for m in out if m.warnings]
    if errs:
        print("\n   ❌ Ingestion errors:")
        for m in errs:
            print(f"      {m.canonical_name} ({Path(m.source_file).name}): {'; '.join(m.errors)}")
    if warns:
        print("\n   ⚠️ Ingestion warnings:")
        for m in warns:
            print(f"      {m.canonical_name}: {'; '.join(m.warnings)}")
    if missing:
        print("\n   ⬜ No report this cycle (heatmap will show a gap):")
        print(f"      {', '.join(missing)}")
    return out


if __name__ == "__main__":
    import sys
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    files = args or [
        "input/RapidImages_mars.pptx",
        "input/Kvaser_Q1.pptx",
    ]
    results = ingest_portfolio(files)
    print("\n--- SUMMARY ---")
    for m in results:
        dr = "yes" if m.deep_read else "no"
        print(f"  {m.canonical_name:18s} kind={m.source_kind:11s} "
              f"text={len(m.body_text):6d}  deep_read={dr}  "
              f"warn={len(m.warnings)} err={len(m.errors)}")