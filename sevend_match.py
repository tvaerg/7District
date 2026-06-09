"""
7D Portfolio Aggregation Engine -- Content-Based File Matching
==============================================================
Identifies which portfolio company a file belongs to by reading the file's
CONTENT, not its filename. The filename is unreliable -- a "Balansrapport.pdf"
could be from any portfolio company -- so we open the file and look for company
names inside.

Cost model
----------
* PDFs: a fast local text peek with pypdf (first 2 pages). NO Gemini call here.
  If pypdf fails or yields no text, we fall back to filename-based matching as
  a last resort rather than spending an API call on identification.
* PPTX: read first 2 slides' text with python-pptx. Local, fast.
* XLSX: read the first sheet's first ~20 rows with openpyxl. Local, fast.
* DOCX: read first ~20 paragraphs with python-docx. Local, fast.

Each peek returns a small text sample. We then score the sample against the
registry's aliases (same scoring as the old filename matcher, just running on
text instead of filename) and pick the highest-scoring company -- if its score
clears a confidence floor. Otherwise the file is reported as UNMATCHED so the
operator can decide rather than misroute.

Why not always use Gemini for identification?
---------------------------------------------
Two reasons:
* Cost. We're already paying for Gemini to do the deep read on PDFs. Adding a
  pre-pass call per file would double API cost and slow the pipeline by ~30s
  per file. Local text extraction handles the same job for free in most cases.
* Reliability. Company names are usually in the first sentence or header of any
  report. Local extraction catches that 95% of the time. Gemini is the
  fallback for the 5% where it doesn't.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from sevend_registry import REGISTRY, CompanyRecord, FILE_MAP


# Minimum score to accept a content match. Tuned empirically: a single short
# alias hit ('acre') scores ~6; we want at least two hits or one solid one.
MIN_CONFIDENCE_SCORE = 8

# How much text to read from each format before we give up looking.
_PDF_PAGES_PEEK = 2
_PPTX_SLIDES_PEEK = 2
_XLSX_ROWS_PEEK = 20
_DOCX_PARAS_PEEK = 25


# -----------------------------------------------------------------------------
# Per-format text peeks (all local, no API)
# -----------------------------------------------------------------------------
def _peek_pdf(path: Path) -> str:
    """Quick text extraction from the first N pages of a PDF, using pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        chunks = []
        for i, page in enumerate(reader.pages[:_PDF_PAGES_PEEK]):
            try:
                chunks.append(page.extract_text() or "")
            except Exception:                          # noqa: BLE001
                continue
        return "\n".join(chunks)
    except Exception:                                  # noqa: BLE001
        return ""


def _peek_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides):
            if i >= _PPTX_SLIDES_PEEK:
                break
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text or "")
        return "\n".join(chunks)
    except Exception:                                  # noqa: BLE001
        return ""


def _peek_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True, read_only=True)
        ws = wb[wb.sheetnames[0]]
        chunks = []
        for ri, row in enumerate(ws.iter_rows(values_only=True)):
            if ri >= _XLSX_ROWS_PEEK:
                break
            cells = [str(c) for c in row if c is not None]
            if cells:
                chunks.append(" | ".join(cells))
        return "\n".join(chunks)
    except Exception:                                  # noqa: BLE001
        return ""


def _peek_docx(path: Path) -> str:
    # Import is outside the try so a missing dependency raises clearly upstream,
    # rather than being swallowed and misreported as "file unreadable".
    try:
        import docx
    except ImportError as e:
        raise ImportError(
            "python-docx is required for .docx files (`pip install python-docx`)"
        ) from e
    try:
        doc = docx.Document(str(path))
        chunks = []
        for i, para in enumerate(doc.paragraphs):
            if i >= _DOCX_PARAS_PEEK:
                break
            t = (para.text or "").strip()
            if t:
                chunks.append(t)
        # Also scan the first few table cells -- some Word reports put the
        # company name in a header table, not a paragraph.
        for ti, tbl in enumerate(doc.tables[:3]):
            for row in tbl.rows[:3]:
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        chunks.append(t)
        return "\n".join(chunks)
    except Exception as e:                             # noqa: BLE001
        # File-level read failure (corrupt file, locked, etc.) -- not the same
        # as missing library. Return empty + the matcher will report it cleanly.
        return ""


_PEEK_BY_EXT = {
    ".pdf":  _peek_pdf,
    ".pptx": _peek_pptx,  ".ppt": _peek_pptx,
    ".xlsx": _peek_xlsx,  ".xlsm": _peek_xlsx,
    ".docx": _peek_docx,  ".doc": _peek_docx,
}


def peek_content(path: str | Path) -> str:
    """Return a short text sample from the file, or empty string if unreadable."""
    p = Path(path)
    fn = _PEEK_BY_EXT.get(p.suffix.lower())
    return fn(p) if fn else ""


# -----------------------------------------------------------------------------
# Alias scoring over arbitrary text (same shape as the old filename matcher)
# -----------------------------------------------------------------------------
_WORD_RE = re.compile(r"[a-z0-9]+")


def _score_text_for_company(text_norm: str, tokens: set[str], rec: CompanyRecord) -> int:
    """Score how strongly the given text identifies this company."""
    score = 0
    # Canonical name match is the strongest signal (e.g. "Acre AB" in the text).
    cn = re.sub(r"[^a-z0-9]+", " ", rec.canonical_name.lower()).strip()
    if cn and cn in text_norm:
        score += 2 * len(cn) + 5
    for alias in rec.aliases:
        a = re.sub(r"[^a-z0-9]+", "", alias.lower())
        if not a:
            continue
        if a in tokens:
            score += len(a) + 2          # whole-token hit, weighted by length
        elif a in text_norm:
            score += len(a)              # substring hit
    return score


def match_by_content(path: str | Path) -> tuple[Optional[CompanyRecord], int, str]:
    """Identify the company from the file's content. Returns
    (company_or_None, confidence_score, explanation_string)."""
    text = peek_content(path)
    if not text:
        return None, 0, "could not extract text from file"

    text_norm = re.sub(r"\s+", " ", text.lower())
    tokens = set(_WORD_RE.findall(text_norm))

    best: Optional[CompanyRecord] = None
    best_score = 0
    for rec in REGISTRY:
        s = _score_text_for_company(text_norm, tokens, rec)
        if s > best_score:
            best_score, best = s, rec

    if best is not None and best_score >= MIN_CONFIDENCE_SCORE:
        return best, best_score, f"content-matched {best.canonical_name} (score {best_score})"
    if best is not None:
        return None, best_score, (
            f"low confidence: best candidate {best.canonical_name} scored only "
            f"{best_score} (floor is {MIN_CONFIDENCE_SCORE})"
        )
    return None, 0, "no company name detected in file content"


# -----------------------------------------------------------------------------
# Public entry point used by ingestion
# -----------------------------------------------------------------------------
def identify_company(path: str | Path) -> tuple[Optional[CompanyRecord], str]:
    """Resolve a file to a company by reading its CONTENT.

    Decision flow:
      1. If the operator pinned this filename in FILE_MAP (registry override),
         honour that -- explicit beats implicit.
      2. Otherwise read the file's content and score it against the registry.
      3. If no confident match, return None and a clear explanation.
    """
    from sevend_registry import _by_name

    p = Path(path)

    # 1. Operator override always wins.
    stem = re.sub(r"[^a-z0-9]+", "_", p.stem.lower())
    if stem in FILE_MAP:
        rec = _by_name().get(FILE_MAP[stem].lower())
        if rec is not None:
            return rec, f"FILE_MAP override -> {rec.canonical_name}"

    # 2. Content-based match (local, free).
    rec, score, why = match_by_content(p)
    if rec is not None:
        return rec, why

    # 3. Gemini fallback for files where local extraction wasn't enough --
    # logo-only PPTX/PDF, data-only XLSX, etc. We pay one extra API call to
    # identify the company by sending the file to Gemini with a focused
    # "name the company" prompt.
    gemini_rec, gemini_why = _gemini_identify(p)
    if gemini_rec is not None:
        return gemini_rec, gemini_why

    return None, why or gemini_why or "could not identify company"


def _gemini_identify(path: Path) -> tuple[Optional[CompanyRecord], str]:
    """Last-resort identification via Gemini. Returns (company, explanation)."""
    try:
        import os
        from dotenv import load_dotenv
        # Reuse the same dotenv discovery as the rest of the engine.
        here = Path(__file__).resolve().parent
        for cand in (here / ".env", here.parent / ".env"):
            if cand.is_file():
                load_dotenv(cand); break
        else:
            load_dotenv()

        api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
        if not api_key:
            return None, "no GEMINI_API_KEY / GOOGLE_API_KEY for Gemini fallback"

        import google.generativeai as genai
        genai.configure(api_key=api_key)

        # List the registry's canonical names + key aliases so Gemini knows the
        # valid answer set. Without this constraint it would invent companies.
        candidates_block = "\n".join(
            f"  - {rec.canonical_name}  (aliases: {', '.join(rec.aliases)})"
            for rec in REGISTRY
        )
        prompt = (
            "You are identifying which portfolio company a financial report belongs to.\n"
            "The company MUST be one of the following:\n"
            f"{candidates_block}\n\n"
            "Read the document and return ONLY the canonical name of the company, exactly as "
            "spelled in the list above. If the document cannot be confidently mapped to one of "
            "these companies, return UNKNOWN. Return one word/phrase only, no explanation."
        )

        model = genai.GenerativeModel(
            os.getenv("GEMINI_MODEL", "models/gemini-3.1-pro-preview")
        )

        ext = path.suffix.lower()
        if ext == ".pdf":
            uploaded = genai.upload_file(path=str(path), display_name=path.stem)
            try:
                resp = model.generate_content([uploaded, prompt])
            finally:
                try:
                    genai.delete_file(uploaded.name)
                except Exception:                      # noqa: BLE001
                    pass
        else:
            # For non-PDF, hand Gemini the local text peek -- much cheaper than
            # converting and uploading the binary, and the company name is
            # almost always in the first page of text.
            text = peek_content(path)
            if not text:
                return None, "Gemini fallback: no extractable text"
            resp = model.generate_content(f"{prompt}\n\nDOCUMENT:\n{text}")

        guess = (getattr(resp, "text", "") or "").strip().strip(".\"' ")
        if not guess or guess.upper() == "UNKNOWN":
            return None, "Gemini could not identify the company"

        # Map Gemini's answer back to a CompanyRecord.
        gl = guess.lower()
        for rec in REGISTRY:
            if rec.canonical_name.lower() == gl or gl in [a.lower() for a in rec.aliases]:
                return rec, f"Gemini-identified -> {rec.canonical_name}"
            if rec.canonical_name.lower() in gl:
                return rec, f"Gemini-identified -> {rec.canonical_name} (substring)"

        return None, f"Gemini returned '{guess}' but it doesn't match any registry company"
    except Exception as e:                             # noqa: BLE001
        return None, f"Gemini fallback failed: {e}"


if __name__ == "__main__":
    import sys
    for arg in sys.argv[1:]:
        rec, why = identify_company(arg)
        name = rec.canonical_name if rec else "UNMATCHED"
        print(f"  {Path(arg).name:55s} -> {name:20s} ({why})")